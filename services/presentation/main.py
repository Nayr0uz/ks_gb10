# services/presentation/main.py

from fastapi import FastAPI, HTTPException
from typing import Optional
from pydantic import BaseModel
import os
import io
from pptx import Presentation
from pptx.util import Inches
from starlette.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from shared.database import get_database
from presentation_service import generate_presentation_content, generate_presentation_content_streaming
import json
import asyncio
from starlette.responses import StreamingResponse
from contextlib import asynccontextmanager
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# ==========================
# DB Manager and Lifespan
# ==========================
db_manager = None

@asynccontextmanager
async def lifespan(app: FastAPI):
    """Application lifespan manager"""
    global db_manager
    logger.info("Starting Presentation Service...")
    
    db_manager = get_database()
    await db_manager.initialize()
    
    logger.info("Presentation Service started successfully")
    yield
    
    logger.info("Shutting down Presentation Service...")
    await db_manager.close()

app = FastAPI(
    title="ENBD Presentation Service",
    description="Service to generate presentations from documents.",
    version="1.0.0",
    lifespan=lifespan
)

# ==========================
# CORS
# ==========================
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ==========================
# expose generated files
# ==========================
os.makedirs("./saved_presentations", exist_ok=True)
app.mount(
    "/saved_presentations",
    StaticFiles(directory="./saved_presentations"),
    name="saved_presentations",
)

# ==========================
# schemas
# ==========================
class PresentationConfig(BaseModel):
    title: str
    scope: str
    topic: Optional[str] = None
    detail_level: str
    difficulty: str
    slide_style: str
    num_slides: int
    include_diagrams: bool
    include_code_examples: bool


# ==========================
# create / generate (SYNC)
# ==========================
@app.post("/api/v1/presentation/generate-presentation-stream/")
async def generate_presentation_stream(config: PresentationConfig):
    
    queue = asyncio.Queue()

    async def slide_callback(slide_index, slide_content):
        slide_data = {"index": slide_index, "content": slide_content}
        await queue.put(f"{json.dumps(slide_data)}\n\n")

    async def event_stream(presentation_id: int):
        # Start generation in the background
        generation_task = asyncio.create_task(generate_presentation_content_streaming(
            presentation_id,
            callback=slide_callback,
        ))
        generation_task.add_done_callback(lambda t: queue.put_nowait(None))
        
        # Yield from queue
        while True:
            data = await queue.get()
            if data is None:
                break
            yield data

    # 1) Save request to Neo4j
    presentation_id = await db_manager.create_presentation(
        title=config.title,
        scope=config.scope,
        topic=config.topic,
        detail_level=config.detail_level,
        difficulty=config.difficulty,
        slide_style=config.slide_style,
        num_slides=config.num_slides,
        include_diagrams=config.include_diagrams,
        include_code_examples=config.include_code_examples
    )

    return StreamingResponse(event_stream(presentation_id), media_type="text/event-stream")



@app.get("/api/v1/presentation/{presentation_id}/download/ppt")
async def download_ppt(presentation_id: str):
    presentation = await db_manager.get_presentation(presentation_id)
    
    if not presentation or not presentation.get("content"):
        raise HTTPException(status_code=404, detail="Presentation or content not found")

    row = {
        "title": presentation.get("title"),
        "content": presentation.get("content")
    }
    if not row or not row["content"]:
        raise HTTPException(status_code=404, detail="Presentation or content not found")

    prs = Presentation()
    slides_content = row["content"].split("\n\n---SLIDE_SEPARATOR---\n\n")

    for slide_text in slides_content:
        # Skip empty or separator-only slides
        if not slide_text.strip() or slide_text.strip() == "---SLIDE_SEPARATOR---":
            continue

        lines = slide_text.strip().split('\n')
        title = lines[0].replace("**", "") if lines else "Slide"
        content = "\n".join(lines[1:]) if len(lines) > 1 else ""

        slide_layout = prs.slide_layouts[5]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        # Safely set title if a title placeholder exists
        try:
            slide.shapes.title.text = title
        except Exception:
            pass

        # Safely set body content. Some layouts may not include the typical
        # body placeholder at index 1, which would raise a KeyError. In that
        # case, add a textbox and insert the content there.
        try:
            slide.placeholders[1].text = content
        except Exception:
            try:
                left = Inches(1)
                top = Inches(1.8)
                width = Inches(8)
                height = Inches(4.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = content
            except Exception:
                # As a last resort, skip adding content for this slide
                # to avoid failing the entire request.
                continue

    f = io.BytesIO()
    prs.save(f)
    f.seek(0)
    
    headers = {
        'Content-Disposition': f'attachment; filename="{row["title"]}.pptx"'
    }

    return StreamingResponse(f, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers=headers)


# ==========================
# get single
# ==========================
@app.get("/api/v1/presentation/{presentation_id}")
async def get_presentation_status(presentation_id: str):
    row = await db_manager.get_presentation(presentation_id)

    if not row:
        raise HTTPException(status_code=404, detail="Presentation not found")

    return {
        "id": row["id"],
        "title": row["title"],
        "status": row["status"],
        "output_file_path": row["output_file_path"],
        "created_at": row["created_at"].isoformat() if row["created_at"] else None,
        "content": row["content"],
    }


# ==========================
# list (for frontend sidebar)
# ==========================
@app.get("/api/v1/presentation/presentations")
async def list_presentations(limit: int = 50):
    rows = await db_manager.list_presentations(limit)

    return [
        {
            "id": r["id"],
            "title": r["title"],
            "status": r["status"],
            "output_file_path": r["output_file_path"],
            "created_at": r["created_at"].isoformat() if r["created_at"] else None,
        }
        for r in rows
    ]


# ==========================
# aliases بدون /api/...
# ==========================
@app.get("/presentations")
async def list_presentations_alt(limit: int = 50):
    return await list_presentations(limit=limit)


@app.get("/presentations/{presentation_id}")
async def get_presentation_status_alt(presentation_id: str):
    return await get_presentation_status(presentation_id)


@app.get("/presentations/{presentation_id}/download/ppt")
async def download_ppt_alt(presentation_id: str):
    return await download_ppt(presentation_id)


# ==========================
# health
# ==========================
@app.get("/api/v1/presentation/health")
async def health_check():
    return {"status": "ok"}


# ==========================
# run
# ==========================
if __name__ == "__main__":
    import uvicorn

    uvicorn.run("main:app", host="0.0.0.0", port=8003, reload=False)
